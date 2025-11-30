"""
Generate PowerPoint presentation for Stock Research Workflow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add dark blue background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0x1a, 0x36, 0x5d)  # Dark blue
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, phase_num=None):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add gradient-like background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0x2d, 0x5a, 0x8a)  # Medium blue
    background.line.fill.background()
    
    # Phase number circle
    if phase_num:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(4), Inches(1.5), Inches(2), Inches(2)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0xff, 0xc1, 0x07)  # Gold
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(4), Inches(1.8), Inches(2), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(phase_num)
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)
        p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, goal=None):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
    background.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1a, 0x36, 0x5d)
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # Goal box if provided
    y_start = Inches(1.5)
    if goal:
        goal_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(0.6)
        )
        goal_box.fill.solid()
        goal_box.fill.fore_color.rgb = RGBColor(0xe8, 0xf4, 0xf8)
        goal_box.line.color.rgb = RGBColor(0x2d, 0x5a, 0x8a)
        
        goal_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(0.5))
        tf = goal_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"🎯 Goal: {goal}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x2d, 0x5a, 0x8a)
        p.font.bold = True
        y_start = Inches(2.2)
    
    # Content items
    content_box = slide.shapes.add_textbox(Inches(0.5), y_start, Inches(9), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(item, dict):
            # Main item with sub-items
            p.text = f"• {item['title']}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1a, 0x36, 0x5d)
            p.space_before = Pt(12)
            
            for sub in item.get('subs', []):
                p = tf.add_paragraph()
                p.text = f"    ◦ {sub}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0x4a, 0x4a, 0x4a)
                p.space_before = Pt(4)
        else:
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            p.space_before = Pt(10)
    
    return slide

def add_checklist_slide(prs, title, items):
    """Add a checklist slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xf0, 0xf4, 0xf8)
    background.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x28, 0xa7, 0x45)  # Green
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"✓ {title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # Checklist items
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"☐ {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(14)
    
    return slide

def add_template_slide(prs, title, steps):
    """Add a research template slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    background.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x6f, 0x42, 0xc1)  # Purple
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📋 {title}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # Steps with numbered boxes
    y_pos = 1.6
    for i, step in enumerate(steps, 1):
        # Number box
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(0.6), Inches(0.6)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = RGBColor(0x6f, 0x42, 0xc1)
        num_box.line.fill.background()
        
        num_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.1), Inches(0.6), Inches(0.5))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(i)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER
        
        # Step text
        step_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos + 0.1), Inches(8), Inches(0.5))
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
        y_pos += 0.9
    
    return slide

def add_overview_slide(prs):
    """Add workflow overview slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
    background.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1a, 0x36, 0x5d)
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "5-Phase Research Process"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # Phase boxes
    phases = [
        ("1", "Discovery", "Find candidates", RGBColor(0x00, 0x7b, 0xff)),
        ("2", "Technical", "Chart analysis", RGBColor(0x28, 0xa7, 0x45)),
        ("3", "Fundamental", "Business quality", RGBColor(0xfd, 0x7e, 0x14)),
        ("4", "SEC Filing", "Verify financials", RGBColor(0xdc, 0x35, 0x45)),
        ("5", "Decision", "Act & monitor", RGBColor(0x6f, 0x42, 0xc1)),
    ]
    
    x_start = 0.3
    box_width = 1.8
    
    for i, (num, title, desc, color) in enumerate(phases):
        x = x_start + i * (box_width + 0.1)
        
        # Phase box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(box_width), Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Number
        num_text = slide.shapes.add_textbox(Inches(x), Inches(1.9), Inches(box_width), Inches(0.8))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_text = slide.shapes.add_textbox(Inches(x), Inches(2.7), Inches(box_width), Inches(0.6))
        tf = title_text.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_text = slide.shapes.add_textbox(Inches(x), Inches(3.3), Inches(box_width), Inches(0.8))
        tf = desc_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow (except last)
        if i < len(phases) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + box_width), Inches(2.8), Inches(0.15), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0x66, 0x66, 0x66)
            arrow.line.fill.background()
    
    # Key tools section
    tools_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
    tf = tools_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Tools: Market News • D/W/M Charts • Multi-TF Gallery • StockCharts • Fundamental Analysis • Business Analysis • Buffett & CANSLIM • 10-K/10-Q Study • Watch List"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    add_title_slide(prs, 
                   "Stock Research Workflow",
                   "A Systematic Approach to Stock Analysis")
    
    # Overview slide
    add_overview_slide(prs)
    
    # Phase 1: Discovery
    add_section_slide(prs, "Discovery & Screening", 1)
    add_content_slide(prs, "Phase 1: Discovery & Screening", [
        {"title": "Start with Market Overview", "subs": [
            "Click Market news for AI-summarized conditions",
            "Check sector rotation and sentiment"
        ]},
        {"title": "Browse Curated Lists", "subs": [
            "Use ◀/▶ to navigate ticker lists",
            "Open D charts for quick visual scan"
        ]},
        {"title": "External Screening", "subs": [
            "Use Finviz Screener via URLs menu",
            "AI Search: 'Find undervalued tech stocks'"
        ]},
    ], goal="Identify potential investment candidates")
    
    # Phase 2: Technical
    add_section_slide(prs, "Technical Analysis", 2)
    add_content_slide(prs, "Phase 2: Technical Analysis", [
        {"title": "Multi-Timeframe Analysis", "subs": [
            "Click Multi-TF for D/W/M view",
            "Look for trend alignment across timeframes"
        ]},
        {"title": "StockCharts Deep Dive", "subs": [
            "SC button for professional charting",
            "Check support/resistance, MAs, volume"
        ]},
        {"title": "Seasonality & Comparison", "subs": [
            "Seasonality tab for historical patterns",
            "Compare button for relative performance"
        ]},
    ], goal="Evaluate price action and chart patterns")
    
    # Phase 3: Fundamental
    add_section_slide(prs, "Fundamental Analysis", 3)
    add_content_slide(prs, "Phase 3: Fundamental Analysis", [
        {"title": "Quick Metrics Review", "subs": [
            "Fundamental Analysis tab",
            "Filter: 'pe ratio market cap revenue'"
        ]},
        {"title": "Business Deep Dive", "subs": [
            "Run BA for AI analysis",
            "Review business model, moat, risks"
        ]},
        {"title": "Investment Framework", "subs": [
            "Buffett & CANSLIM tab",
            "Check radar chart for strengths/weaknesses"
        ]},
    ], goal="Understand the business and valuation")
    
    # Phase 4: SEC
    add_section_slide(prs, "SEC Filing Analysis", 4)
    add_content_slide(prs, "Phase 4: SEC Filing Analysis", [
        {"title": "Annual Report (10-K)", "subs": [
            "10K Study for AI analysis",
            "Extract Tables for detailed data"
        ]},
        {"title": "Quarterly Report (10-Q)", "subs": [
            "10-Q Study for recent quarter",
            "Compare to previous quarters"
        ]},
        {"title": "Financial Tables", "subs": [
            "Export to Excel for analysis",
            "Track key metrics over time"
        ]},
    ], goal="Verify financials and identify risks")
    
    # Phase 5: Decision
    add_section_slide(prs, "Decision & Monitoring", 5)
    add_content_slide(prs, "Phase 5: Decision & Monitoring", [
        {"title": "Build Your Thesis", "subs": [
            "Summarize findings with clipboard AI",
            "Document: Why buy? At what price? Risk?"
        ]},
        {"title": "Add to Watch List", "subs": [
            "Right-click → Copy to Watch List",
            "Persists across sessions"
        ]},
        {"title": "Ongoing Monitoring", "subs": [
            "Check Stock news regularly",
            "Re-run analysis after earnings"
        ]},
    ], goal="Make informed decision and track position")
    
    # Checklists
    add_section_slide(prs, "Quick Checklists", None)
    
    add_checklist_slide(prs, "5-Minute Stock Check", [
        "Open D chart - Trending up/down/sideways?",
        "Fundamental Analysis - P/E reasonable? Growing revenue?",
        "Stock news - Any red flags?"
    ])
    
    add_checklist_slide(prs, "30-Minute Deep Dive", [
        "Multi-TF charts - Trend alignment?",
        "Seasonality - Good entry timing?",
        "Run BA - Business quality?",
        "Buffett & CANSLIM - Investment grade?",
        "10-Q Study - Recent quarter healthy?"
    ])
    
    add_checklist_slide(prs, "Pre-Earnings Checklist", [
        "Review last 4 quarters via 10-Q Study",
        "Check news for analyst expectations",
        "Review Seasonality for historical reactions",
        "Set position size based on volatility"
    ])
    
    # Templates
    add_section_slide(prs, "Research Templates", None)
    
    add_template_slide(prs, "Value Investing Template", [
        "Fundamental Analysis → Filter: 'book value debt equity roe'",
        "Buffett & CANSLIM → Check value metrics",
        "10-K Study → Verify balance sheet strength",
        "AI Search → 'What are the competitive advantages?'"
    ])
    
    add_template_slide(prs, "Growth Investing Template", [
        "Fundamental Analysis → Filter: 'revenue growth earnings'",
        "Multi-TF Charts → Confirm uptrend",
        "Run BA → Analyze market opportunity",
        "10-Q Study → Verify growth acceleration"
    ])
    
    add_template_slide(prs, "Dividend Investing Template", [
        "Fundamental Analysis → Filter: 'dividend yield payout'",
        "10-K Study → Check dividend history",
        "AI Search → 'Is dividend sustainable?'",
        "Seasonality → Best entry points historically"
    ])
    
    # Final slide
    add_title_slide(prs, 
                   "Start Researching!",
                   "Launch main.py and follow the workflow")
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), "Stock_Research_Workflow.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    main()
