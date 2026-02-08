"""
Convert PITCH_DECK.md to PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

# Color scheme
PRIMARY_BLUE = RGBColor(0x1a, 0x73, 0xe8)  # Google Blue
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xff, 0xff, 0xff)
ACCENT_GREEN = RGBColor(0x34, 0xa8, 0x53)  # Success green


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_items, has_table=False):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    first = True
    for item in content_items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        # Handle different item types
        if item.startswith("###"):
            p.text = item.replace("###", "").strip()
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_BLUE
            p.space_before = Pt(12)
        elif item.startswith("-") or item.startswith("•"):
            p.text = "• " + item.lstrip("-•").strip()
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.level = 0
        elif item.startswith(">"):
            p.text = item.lstrip(">").strip()
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = LIGHT_GRAY
        else:
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
    
    return slide


def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    first = True
    for item in left_items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = "• " + item if not item.startswith("•") else item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    first = True
    for item in right_items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = "• " + item if not item.startswith("•") else item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    return slide


def create_pitch_deck():
    """Create the full pitch deck presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "📈 Personal AI Stock Assistant",
        "Your AI-Powered Investment Research Companion\n\nDec 2025"
    )
    
    # Slide 2: The Problem
    add_content_slide(prs, "😫 The Problem", [
        "### Individual investors face overwhelming challenges:",
        "",
        "- Information Overload: 1000s of stocks, endless news, complex filings",
        "- Time-Consuming Research: Hours spent on each stock analysis",
        "- Scattered Tools: Switching between 10+ websites and apps",
        "- SEC Filing Complexity: 10-K/10-Q reports are 100+ pages",
        "- No AI Integration: Manual analysis without AI assistance",
        "",
        "> Result: Retail investors miss opportunities or make uninformed decisions"
    ])
    
    # Slide 3: The Solution
    add_content_slide(prs, "🚀 The Solution", [
        "### Personal AI Stock Assistant - An all-in-one desktop application:",
        "",
        "- 📊 Multi-timeframe charting (Daily/Weekly/Monthly)",
        "- 🤖 AI-powered analysis using Google Gemini",
        "- 📑 SEC filing extraction (10-K, 10-Q automated parsing)",
        "- 📰 News aggregation & summarization",
        "- 🎯 Investment frameworks (Buffett & CANSLIM analysis)",
        "- 💼 Portfolio management with watch lists",
        "",
        "> One tool. Complete research workflow. AI-enhanced insights."
    ])
    
    # Slide 4: Key Features
    add_two_column_slide(prs, "✨ Key Features", [
        "Smart Charting:",
        "  Daily, Weekly, Monthly candlesticks",
        "  Multi-timeframe galleries",
        "  StockCharts.com integration",
        "  Seasonality analysis",
        "",
        "AI-Powered Analysis:",
        "  Business analysis with Gemini AI",
        "  News summarization (4 types)",
        "  Custom AI search queries",
        "  Clipboard content analysis",
    ], [
        "SEC Filing Intelligence:",
        "  Automated 10-K/10-Q extraction",
        "  Financial table parsing",
        "  AI-generated filing summaries",
        "  Export to Excel",
        "",
        "Investment Frameworks:",
        "  Warren Buffett value metrics",
        "  CANSLIM growth criteria",
        "  Radar chart visualization",
        "  Fundamental data filtering",
    ])
    
    # Slide 5: How It Works
    add_content_slide(prs, "⚙️ 5-Phase Research Workflow", [
        "### Systematic approach to stock analysis:",
        "",
        "1️⃣ DISCOVER → Screen & filter potential investments",
        "2️⃣ TECHNICAL → Chart analysis across timeframes",
        "3️⃣ FUNDAMENTAL → Business & financial analysis",
        "4️⃣ SEC DEEP DIVE → 10-K/10-Q filing analysis",
        "5️⃣ MONITOR → Track positions & make decisions",
        "",
        "### Time Savings:",
        "- 5-minute quick check → Instant stock overview",
        "- 30-minute deep dive → Complete analysis",
        "- Pre-earnings prep → Automated quarterly review"
    ])
    
    # Slide 6: Technology Stack
    add_content_slide(prs, "🛠️ Technology Stack", [
        "### Core Technologies:",
        "",
        "- Frontend: Python Tkinter (Desktop GUI)",
        "- AI Engine: Google Gemini API",
        "- Data Sources: Yahoo Finance, Finviz, SEC EDGAR",
        "- Charting: Plotly, Matplotlib, StockCharts",
        "- Data Processing: Pandas, BeautifulSoup",
        "",
        "### Key Technical Features:",
        "- ✅ Offline-capable with local caching",
        "- ✅ Rate limiting for API compliance",
        "- ✅ Multi-language support (English/Chinese)",
        "- ✅ Persistent settings & watch lists",
        "- ✅ 80+ tooltips for user guidance"
    ])
    
    # Slide 7: Market Opportunity
    add_content_slide(prs, "📈 Market Opportunity", [
        "### Target Market: Self-Directed Investors",
        "",
        "- US Retail Investors: 150M+ accounts (Growing 15% YoY)",
        "- Active Traders: 30M+ (Post-pandemic surge)",
        "- AI Tool Adoption: 40% interested (Rapidly increasing)",
        "",
        "### Why Now?",
        "- 🚀 AI capabilities have reached practical utility",
        "- 📱 Retail trading at all-time highs",
        "- 📊 Demand for sophisticated tools at consumer prices",
        "- 🏦 Democratization of institutional-grade analysis"
    ])
    
    # Slide 8: Competitive Advantage
    add_content_slide(prs, "🏆 Competitive Advantage", [
        "### vs Bloomberg ($24K/yr), Yahoo Finance, TradingView:",
        "",
        "- ✅ AI Analysis - Built-in Gemini integration",
        "- ✅ SEC Parsing - Automated filing extraction",
        "- ✅ News Summary - AI-powered, not manual",
        "- ✅ Buffett/CANSLIM - Investment frameworks built-in",
        "- ✅ Desktop App - Works offline",
        "- ✅ Affordable - Free/low-cost tiers",
        "",
        "### Our Moat:",
        "- AI-First Design - Built around Gemini integration",
        "- SEC Intelligence - Automated filing analysis",
        "- Investment Frameworks - Proven methodologies",
        "- All-in-One - No tab switching between tools"
    ])
    
    # Slide 9: Business Model
    add_content_slide(prs, "💰 Business Model", [
        "### Freemium SaaS Model:",
        "",
        "- Free ($0): Basic charting, limited AI queries",
        "- Pro ($19/mo): Unlimited AI, SEC analysis, priority support",
        "- Team ($49/mo): Multi-user, shared watchlists, API access",
        "",
        "### Revenue Streams:",
        "- Subscriptions: Primary revenue (80%)",
        "- API Access: Developer integrations (15%)",
        "- Premium Data: Enhanced datasets (5%)",
        "",
        "### Unit Economics:",
        "- CAC: ~$30 | LTV: ~$300 | LTV/CAC: 10x"
    ])
    
    # Slide 10: Traction
    add_content_slide(prs, "📊 Traction & Milestones", [
        "### Current Status: MVP Complete ✅",
        "",
        "- Features Built: 50+ core features",
        "- Code Quality: 7,000+ lines, production-ready",
        "- Documentation: Complete user guide",
        "- AI Integration: Gemini API fully integrated",
        "- SEC Parsing: 10-K/10-Q extraction working",
        "",
        "### Roadmap:",
        "- Q1 2025: Beta Launch → 100 users",
        "- Q2 2025: Public Launch → 1,000 users",
        "- Q3 2025: Mobile App → 10,000 users",
        "- Q4 2025: Enterprise → B2B sales"
    ])
    
    # Slide 11: Team
    add_content_slide(prs, "👥 Team", [
        "### Founder:",
        "",
        "Jue Shi",
        "- Software Engineer with finance passion",
        "- Built end-to-end from concept to MVP",
        "- Deep understanding of retail investor needs",
        "",
        "### Advisors Needed:",
        "- 🎯 Finance/Trading expert",
        "- 📈 Growth marketing specialist",
        "- 🏢 Enterprise sales leader"
    ])
    
    # Slide 12: The Ask
    add_content_slide(prs, "🎯 The Ask", [
        "### Seeking: $250K Seed Round",
        "",
        "### Use of Funds:",
        "- Engineering (50%): Web/mobile development",
        "- Marketing (25%): User acquisition",
        "- Infrastructure (15%): Cloud, APIs, data",
        "- Operations (10%): Legal, admin",
        "",
        "### Milestones This Funding Enables:",
        "- ✅ Web application launch",
        "- ✅ 1,000 paying subscribers",
        "- ✅ Mobile app MVP",
        "- ✅ Series A readiness"
    ])
    
    # Slide 13: Vision
    add_content_slide(prs, "🌟 Vision", [
        "> \"Democratize institutional-grade investment research through AI\"",
        "",
        "### Long-term Goals:",
        "- 2025: Best AI stock research tool for individuals",
        "- 2026: Expand to crypto, forex, commodities",
        "- 2027: Enterprise offering for RIAs and advisors",
        "- 2028: Global expansion (EU, Asia markets)",
        "",
        "### Impact:",
        "- 🎯 Help 1M+ investors make better decisions",
        "- 💡 Reduce information asymmetry",
        "- 🌍 Level the playing field vs institutions"
    ])
    
    # Slide 14: Contact
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_BLUE
    shape.line.fill.background()
    
    # Thank you
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(2))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📈 Personal AI Stock Assistant"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nJue Shi"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "\"Empowering individual investors with AI-powered research\""
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = "PITCH_DECK.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_pitch_deck()
